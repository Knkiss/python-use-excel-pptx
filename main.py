from pptx import Presentation
from openpyxl import load_workbook
from pptx.dml.color import RGBColor
import random


def createNewPPT(titleStr, textStr, R, G, B):

    if R is not int:
        R = int(R, 16)
    if G is not int:
        G = int(G, 16)
    if B is not int:
        B = int(B, 16)

    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = RGBColor(R, G, B)
    slide1.shapes.title.text = titleStr
    if sheet[text].value is not None:
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.background.fill.solid()
        slide1.background.fill.fore_color.rgb = RGBColor(R, G, B)
        slide1.shapes.title.text = titleStr
        slide1.placeholders[1].text = textStr


if __name__ == '__main__':
    wb = load_workbook(filename="test1.xlsx")
    sheet = wb["Sheet1"]
    prs = Presentation()
    color = [
        ('F1', 'FA', 'FA'),
        ('E8', 'FF', 'E8'),
        ('E8', 'E8', 'FF'),
        ('F2', 'F1', 'D7'),
        ('FB', 'FB', 'EA'),
        ('D5', 'F3', 'F4')
    ]
    for i in range(1, 94):
        title = 'A' + str(i)
        text = 'B' + str(i)
        createNewPPT(sheet[title].value, sheet[text].value, *random.choice(color))

    prs.save('test1.pptx')
    wb.close()
